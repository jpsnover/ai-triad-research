#!/usr/bin/env python3
"""Generate the AI Triad pitch deck as a .pptx file."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# -- Palette --
BG_DARK   = RGBColor(0x1A, 0x1A, 0x2E)
BG_MED    = RGBColor(0x16, 0x21, 0x3E)
BG_LIGHT  = RGBColor(0x1F, 0x2B, 0x4D)
ACCENT    = RGBColor(0x00, 0xD4, 0xAA)  # teal
ACCENT2   = RGBColor(0xFF, 0x6B, 0x6B)  # coral
ACCENT3   = RGBColor(0x4E, 0xCF, 0xE0)  # sky blue
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY      = RGBColor(0xAA, 0xAA, 0xBB)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xDD)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_fill(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        from lxml import etree
        nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
        solid = shape.fill._fill.find('.//a:solidFill', nsmap)
        if solid is not None:
            clr = solid[0]
            alpha_elem = etree.SubElement(clr, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
            alpha_elem.set('val', str(int(alpha * 1000)))
    return shape


def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def set_text(tf, text, size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return p


def add_para(tf, text, size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
             space_before=Pt(6), space_after=Pt(2), font_name="Calibri", level=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    p.level = level
    return p


def add_bullet(tf, text, size=16, color=WHITE, bold=False, level=0):
    return add_para(tf, text, size=size, color=color, bold=bold, level=level,
                    space_before=Pt(4), space_after=Pt(2))


def section_label(slide, text):
    """Small section label in top-left."""
    tb = add_textbox(slide, Inches(0.6), Inches(0.3), Inches(3), Inches(0.4))
    set_text(tb.text_frame, text, size=11, color=ACCENT, bold=True)


# ============================================================
# SLIDE 1 — COVER
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, BG_DARK)

# Accent bar top
add_shape_fill(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT)

# Title
tb = add_textbox(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.5))
set_text(tb.text_frame, "AI TRIAD", size=60, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Subtitle
tb = add_textbox(slide, Inches(2), Inches(3.3), Inches(9), Inches(0.8))
set_text(tb.text_frame, "Rosetta Stone for AI Policy Discourse", size=28, color=ACCENT, alignment=PP_ALIGN.CENTER)

# Tagline
tb = add_textbox(slide, Inches(2.5), Inches(4.3), Inches(8), Inches(0.7))
set_text(tb.text_frame, "We are not engineering the solution;\nwe are engineering the table at which the solution can be negotiated.",
         size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

# Author
tb = add_textbox(slide, Inches(3), Inches(5.8), Inches(7), Inches(0.9))
tf = tb.text_frame
set_text(tf, "Jeffrey Snover", size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_para(tf, "Inventor of PowerShell  |  Former Technical Fellow, Microsoft  |  Former DE, Google",
         size=13, color=GRAY, alignment=PP_ALIGN.CENTER)
add_para(tf, "@jsnover  |  Berkman Klein Center Fellow",
         size=13, color=GRAY, alignment=PP_ALIGN.CENTER)

# Bottom accent bar
add_shape_fill(slide, Inches(0), Inches(7.44), SLIDE_W, Inches(0.06), ACCENT)


# ============================================================
# SLIDE 2 — MARKET OVERVIEW
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
section_label(slide, "INTROS & STATUS QUO")

tb = add_textbox(slide, Inches(0.8), Inches(0.8), Inches(11.5), Inches(0.7))
set_text(tb.text_frame, "The Market Is Moving Fast. The Conversation Isn't.", size=32, color=WHITE, bold=True)

# Three columns of facts
col_data = [
    ("Policy Acceleration", ACCENT,
     ["EU AI Act enforced Aug 2025",
      "US executive orders on AI (2023-2025)",
      "GPAI Code of Practice (2025)",
      "China, UK, Brazil, India — all regulating simultaneously"]),
    ("Research Explosion", ACCENT3,
     ["Hundreds of policy papers per quarter",
      "Think tanks, labs, governments all publishing",
      "Contradictory conclusions from shared data",
      "No mechanism to reconcile across worldviews"]),
    ("Taxonomy Gap", ACCENT2,
     ["NIST AI RMF — risk-centric only",
      "EU AI Act — risk-centric only",
      "MIT AI Risk Repository — risk-centric only",
      "None capture Accelerationist goals or Skeptic equity concerns"]),
]

for i, (title, accent_c, bullets) in enumerate(col_data):
    left = Inches(0.8 + i * 4.1)
    # Card background
    add_shape_fill(slide, left, Inches(1.9), Inches(3.7), Inches(4.5), BG_LIGHT)
    # Accent top bar on card
    add_shape_fill(slide, left, Inches(1.9), Inches(3.7), Inches(0.06), accent_c)

    tb = add_textbox(slide, left + Inches(0.3), Inches(2.15), Inches(3.1), Inches(0.5))
    set_text(tb.text_frame, title, size=18, color=accent_c, bold=True)

    tb = add_textbox(slide, left + Inches(0.3), Inches(2.7), Inches(3.1), Inches(3.5))
    tf = tb.text_frame
    tf.word_wrap = True
    for j, b in enumerate(bullets):
        if j == 0:
            set_text(tf, f"  {b}", size=14, color=LIGHT_GRAY)
        else:
            add_bullet(tf, f"  {b}", size=14, color=LIGHT_GRAY)

# Bottom callout
tb = add_textbox(slide, Inches(1.5), Inches(6.6), Inches(10), Inches(0.6))
set_text(tb.text_frame,
         "Stakeholders must synthesize across incompatible worldviews — and no tooling exists to help.",
         size=16, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 3 — THE PROBLEM
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
section_label(slide, "INTROS & STATUS QUO")

tb = add_textbox(slide, Inches(0.8), Inches(0.8), Inches(11.5), Inches(0.7))
set_text(tb.text_frame, "The Problem: Three Monologues, Zero Dialogue", size=32, color=WHITE, bold=True)

# Three POV boxes
povs = [
    ("ACCELERATIONIST", "acc", ACCENT,
     "AI is a revolutionary force for solving\nexistential challenges — climate, disease,\nhuman capability."),
    ("SAFETYIST", "saf", ACCENT2,
     "AI poses catastrophic, irreversible,\npotentially existential risks and\npotential loss of control."),
    ("SKEPTIC", "skp", ACCENT3,
     "Future-gazing distracts from immediate\nharms: bias, labor displacement,\nprivacy erosion, over-hype."),
]

for i, (label, prefix, color, desc) in enumerate(povs):
    left = Inches(0.8 + i * 4.1)
    add_shape_fill(slide, left, Inches(1.8), Inches(3.7), Inches(2.2), BG_LIGHT)
    add_shape_fill(slide, left, Inches(1.8), Inches(0.08), Inches(2.2), color)  # left accent

    tb = add_textbox(slide, left + Inches(0.3), Inches(1.95), Inches(3.2), Inches(0.4))
    set_text(tb.text_frame, label, size=16, color=color, bold=True)

    tb = add_textbox(slide, left + Inches(0.3), Inches(2.4), Inches(3.2), Inches(1.5))
    set_text(tb.text_frame, desc, size=13, color=LIGHT_GRAY)

# The key insight box
add_shape_fill(slide, Inches(0.8), Inches(4.3), Inches(11.7), Inches(2.6), BG_MED)

tb = add_textbox(slide, Inches(1.2), Inches(4.5), Inches(11), Inches(2.4))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "They use the same words — \"harm,\" \"risk,\" \"alignment\" — to mean entirely different things.",
         size=18, color=WHITE, bold=True)
add_para(tf, "", size=8, color=WHITE)
add_para(tf, "A policymaker today cannot distinguish whether a disagreement is rooted in:",
         size=15, color=LIGHT_GRAY)
add_bullet(tf, "Values/Goals  —  fundamentally different desired end-states", size=15, color=LIGHT_GRAY, level=1)
add_bullet(tf, "Data/Facts  —  different beliefs about what is objectively true", size=15, color=LIGHT_GRAY, level=1)
add_bullet(tf, "Method  —  different ways of interpreting the same data", size=15, color=LIGHT_GRAY, level=1)
add_para(tf, "", size=8, color=WHITE)
add_para(tf, "Cost: Weeks of analyst time per question. Blind spots invisible until they surface in testimony or board memos.",
         size=14, color=ACCENT2, bold=True)


# ============================================================
# SLIDE 4 — THE SOLUTION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
section_label(slide, "PRODUCT")

tb = add_textbox(slide, Inches(0.8), Inches(0.8), Inches(11.5), Inches(0.7))
set_text(tb.text_frame, "The Solution: A Translation Layer for AI Policy", size=32, color=WHITE, bold=True)

tb = add_textbox(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(0.6))
set_text(tb.text_frame,
         "AI Triad maps three distinct moral vocabularies onto a single unified schema.\nIt treats conflict as data, not failure.",
         size=16, color=GRAY)

# Three-step pipeline
steps = [
    ("1", "INGEST & STRUCTURE", ACCENT,
     "Converts any source document\n(PDF, HTML, DOCX) into structured\nmarkdown. Extracts metadata,\nPOV tags, topic classification."),
    ("2", "ANALYZE PER POV", ACCENT3,
     "AI generates per-perspective\nsummaries with factual claims,\ntemporal scoping, and AIF-aligned\nargumentation mapping."),
    ("3", "DETECT CONFLICTS", ACCENT2,
     "Automatically surfaces where\nsources disagree and classifies\nthe root: Values vs. Facts vs.\nMethod. Claim-level precision."),
]

for i, (num, title, color, desc) in enumerate(steps):
    left = Inches(0.8 + i * 4.1)
    add_shape_fill(slide, left, Inches(2.6), Inches(3.7), Inches(3.6), BG_LIGHT)

    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(1.5), Inches(2.85), Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    ctf = circle.text_frame
    ctf.paragraphs[0].text = num
    ctf.paragraphs[0].font.size = Pt(24)
    ctf.paragraphs[0].font.bold = True
    ctf.paragraphs[0].font.color.rgb = BG_DARK
    ctf.paragraphs[0].alignment = PP_ALIGN.CENTER
    ctf.vertical_anchor = MSO_ANCHOR.MIDDLE

    tb = add_textbox(slide, left + Inches(0.3), Inches(3.7), Inches(3.1), Inches(0.4))
    set_text(tb.text_frame, title, size=15, color=color, bold=True, alignment=PP_ALIGN.CENTER)

    tb = add_textbox(slide, left + Inches(0.3), Inches(4.2), Inches(3.1), Inches(1.8))
    set_text(tb.text_frame, desc, size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Arrows between steps
for i in range(2):
    left = Inches(4.5 + i * 4.1)
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, Inches(4.1), Inches(0.5), Inches(0.3))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = GRAY
    arrow.line.fill.background()

# Bottom output
tb = add_textbox(slide, Inches(1.5), Inches(6.5), Inches(10), Inches(0.7))
set_text(tb.text_frame,
         "Output: A researcher sees the full argumentative landscape for any topic in minutes, not weeks.",
         size=16, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 5 — PRODUCT SHOWCASE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
section_label(slide, "PRODUCT")

tb = add_textbox(slide, Inches(0.8), Inches(0.8), Inches(11.5), Inches(0.7))
set_text(tb.text_frame, "Product Showcase", size=32, color=WHITE, bold=True)

features = [
    ("Taxonomy Editor", ACCENT,
     ["Desktop app (Electron + React)",
      "Browse, edit, debate across four-POV taxonomy",
      "Edge browser: 7 canonical AIF-aligned relationship types",
      "BDI-structured debate agents (Beliefs, Desires, Intentions)",
      "Graph attribute panels, cross-cutting conflict views"]),
    ("Document Pipeline", ACCENT3,
     ["Fully scriptable CLI (PowerShell 7)",
      "Import -> Summarize -> Detect Conflicts",
      "Multi-AI backend: Gemini, Claude, Groq",
      "Per-POV summaries with factual claim extraction",
      "Temporal scoping on all claims"]),
    ("Conflict Analysis", ACCENT2,
     ["Side-by-side claim comparison across POVs",
      "Attack-type classification: rebut / undercut / undermine",
      "Disagreement tagging: definitional / interpretive / structural",
      "Root-cause diagnosis: Values vs. Facts vs. Method",
      "Compounding data asset — richer with every document"]),
    ("Rosetta Lens (Stretch)", RGBColor(0xBB, 0x86, 0xFC),
     ["Consensus heatmap overlaid on source text",
      "Color-codes by cross-POV alignment level",
      "Perspective pivoting: select text, see other factions' views",
      "Visual diagnostic for policymakers and press",
      ""]),
]

for i, (title, color, bullets) in enumerate(features):
    col = i % 2
    row = i // 2
    left = Inches(0.8 + col * 6.2)
    top = Inches(1.7 + row * 2.7)
    add_shape_fill(slide, left, top, Inches(5.8), Inches(2.5), BG_LIGHT)
    add_shape_fill(slide, left, top, Inches(5.8), Inches(0.05), color)

    tb = add_textbox(slide, left + Inches(0.3), top + Inches(0.15), Inches(5.2), Inches(0.4))
    set_text(tb.text_frame, title, size=16, color=color, bold=True)

    tb = add_textbox(slide, left + Inches(0.3), top + Inches(0.55), Inches(5.2), Inches(1.9))
    tf = tb.text_frame
    tf.word_wrap = True
    for j, b in enumerate(bullets):
        if not b:
            continue
        if j == 0:
            set_text(tf, f"  {b}", size=12, color=LIGHT_GRAY)
        else:
            add_bullet(tf, f"  {b}", size=12, color=LIGHT_GRAY)


# ============================================================
# SLIDE 6 — TARGET AUDIENCE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
section_label(slide, "PRODUCT")

tb = add_textbox(slide, Inches(0.8), Inches(0.8), Inches(11.5), Inches(0.7))
set_text(tb.text_frame, "Target Audience", size=32, color=WHITE, bold=True)

audiences = [
    ("PRIMARY", "Policy Analysts", ACCENT,
     "AI policy analysts at think tanks, government agencies\n(e.g., congressional research staff), and corporate\nAI governance teams who must brief decision-makers\non contested topics.",
     "Gain a diagnostic tool to parse arguments — distinguish\nwhether friction is rooted in values, facts, or method."),
    ("SECONDARY", "Academic Researchers", ACCENT3,
     "Researchers in STS, AI ethics, and technology policy\nwho need systematic multi-perspective literature\ncomparison across hundreds of sources.",
     "Systematic cross-POV analysis replaces weeks of\nmanual reading with structured, claim-level comparison."),
    ("TERTIARY", "AI Journalists", ACCENT2,
     "Journalists covering AI topics who need to quickly\nparse competing narratives and identify where\nsources are genuinely arguing past each other.",
     "A \"decoder ring\" to spot when value conflicts are\ndisguised as factual disputes — report on real trade-offs,\nnot amplify incompatible monologues."),
]

for i, (tier, name, color, desc, value) in enumerate(audiences):
    top = Inches(1.7 + i * 1.8)
    add_shape_fill(slide, Inches(0.8), top, Inches(11.7), Inches(1.6), BG_LIGHT)
    add_shape_fill(slide, Inches(0.8), top, Inches(0.08), Inches(1.6), color)

    # Tier label
    tb = add_textbox(slide, Inches(1.1), top + Inches(0.1), Inches(1.2), Inches(0.3))
    set_text(tb.text_frame, tier, size=10, color=color, bold=True)

    # Name
    tb = add_textbox(slide, Inches(1.1), top + Inches(0.35), Inches(2), Inches(0.4))
    set_text(tb.text_frame, name, size=18, color=WHITE, bold=True)

    # Description
    tb = add_textbox(slide, Inches(3.3), top + Inches(0.15), Inches(4.5), Inches(1.3))
    set_text(tb.text_frame, desc, size=12, color=LIGHT_GRAY)

    # Value prop
    tb = add_textbox(slide, Inches(8), top + Inches(0.15), Inches(4.3), Inches(1.3))
    set_text(tb.text_frame, value, size=12, color=color)


# ============================================================
# SLIDE 7 — ROADMAP
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
section_label(slide, "PRODUCT")

tb = add_textbox(slide, Inches(0.8), Inches(0.8), Inches(11.5), Inches(0.7))
set_text(tb.text_frame, "Roadmap", size=32, color=WHITE, bold=True)

# Horizontal timeline
phases = [
    ("NOW", ACCENT, [
        "Desktop app + CLI pipeline",
        "Four-POV taxonomy (acc/saf/skp/cc)",
        "Automated conflict detection",
        "Debate synthesis with AIF argument mapping",
        "BDI-structured agent debates",
        "Multi-AI backend (Gemini, Claude, Groq)",
    ]),
    ("NEXT", ACCENT3, [
        "Merge POViewer & Summary Viewer into Taxonomy Editor",
        "Rosetta Lens visualization prototype",
        "Advocacy campaign (NeurIPS, policy editorials)",
        "BKC AI Triad Working Group launch",
        "2026 AI Triad Rosetta Stone paper",
        "",
    ]),
    ("LATER", RGBColor(0xBB, 0x86, 0xFC), [
        "Web version for broader access",
        "Personal POVs — define your own perspective",
        "Retarget to other policy domains (energy, biotech)",
        "Real-time monitoring of new publications",
        "API access for institutional integration",
        "",
    ]),
]

# Timeline bar
add_shape_fill(slide, Inches(0.8), Inches(2.1), Inches(11.7), Inches(0.06), GRAY)

for i, (label, color, items) in enumerate(phases):
    left = Inches(0.8 + i * 4.1)

    # Circle on timeline
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(1.65), Inches(1.9), Inches(0.4), Inches(0.4))
    c.fill.solid()
    c.fill.fore_color.rgb = color
    c.line.fill.background()

    # Label
    tb = add_textbox(slide, left, Inches(2.4), Inches(3.7), Inches(0.5))
    set_text(tb.text_frame, label, size=22, color=color, bold=True, alignment=PP_ALIGN.CENTER)

    # Card
    add_shape_fill(slide, left, Inches(3.0), Inches(3.7), Inches(3.8), BG_LIGHT)
    add_shape_fill(slide, left, Inches(3.0), Inches(3.7), Inches(0.05), color)

    tb = add_textbox(slide, left + Inches(0.3), Inches(3.2), Inches(3.1), Inches(3.5))
    tf = tb.text_frame
    tf.word_wrap = True
    for j, item in enumerate(items):
        if not item:
            continue
        if j == 0:
            set_text(tf, f"  {item}", size=13, color=LIGHT_GRAY)
        else:
            add_bullet(tf, f"  {item}", size=13, color=LIGHT_GRAY)


# ============================================================
# SLIDE 8 — COMPETITOR DYNAMICS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
section_label(slide, "WHY US")

tb = add_textbox(slide, Inches(0.8), Inches(0.8), Inches(11.5), Inches(0.7))
set_text(tb.text_frame, "Competitive Landscape", size=32, color=WHITE, bold=True)

# Table header
headers = ["Capability", "Generic AI Summarizers\n(ChatGPT, NotebookLM)", "Policy Databases\n(GovTrack, PolicyNet)",
           "Risk Taxonomies\n(NIST RMF, EU AI Act)", "AI Triad"]
rows_data = [
    ["Multi-POV structure", "No — single\n\"neutral\" summary", "No", "Risk-centric only", "Yes — four named\nperspectives"],
    ["Cross-source conflict\ndetection", "No", "No", "No", "Yes — automated,\nclaim-level"],
    ["Root-cause diagnosis\n(values/facts/method)", "No", "No", "No", "Yes"],
    ["Argumentation\nmapping (AIF)", "No", "No", "No", "Yes — typed edges,\nattack classification"],
]

col_widths = [Inches(2.2), Inches(2.4), Inches(2.2), Inches(2.4), Inches(2.4)]
row_h = Inches(0.85)
start_left = Inches(0.8)
start_top = Inches(1.8)

# Header row
x = start_left
for ci, (hdr, w) in enumerate(zip(headers, col_widths)):
    bg_c = ACCENT if ci == 4 else BG_MED
    add_shape_fill(slide, x, start_top, w, Inches(0.7), bg_c)
    tb = add_textbox(slide, x + Inches(0.1), start_top + Inches(0.05), w - Inches(0.2), Inches(0.6))
    txt_c = BG_DARK if ci == 4 else WHITE
    set_text(tb.text_frame, hdr, size=11, color=txt_c, bold=True, alignment=PP_ALIGN.CENTER)
    x += w

# Data rows
for ri, row in enumerate(rows_data):
    y = start_top + Inches(0.7) + ri * row_h
    x = start_left
    for ci, (cell, w) in enumerate(zip(row, col_widths)):
        bg_c = BG_LIGHT if ri % 2 == 0 else BG_MED
        if ci == 4:
            bg_c = RGBColor(0x0A, 0x3D, 0x32) if ri % 2 == 0 else RGBColor(0x0C, 0x4A, 0x3C)
        add_shape_fill(slide, x, y, w, row_h, bg_c)
        tb = add_textbox(slide, x + Inches(0.1), y + Inches(0.05), w - Inches(0.2), row_h - Inches(0.1))
        txt_c = ACCENT if ci == 4 and cell.startswith("Yes") else (ACCENT2 if cell == "No" else LIGHT_GRAY)
        if ci == 0:
            txt_c = WHITE
        set_text(tb.text_frame, cell, size=11, color=txt_c, bold=(ci == 4 and cell.startswith("Yes")),
                 alignment=PP_ALIGN.CENTER)
        x += w

# Bottom insight
tb = add_textbox(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.8))
set_text(tb.text_frame,
         "Every existing framework is risk-centric. None capture the full Accelerationist-Safetyist-Skeptic landscape.",
         size=16, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 9 — SECRET SAUCE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
section_label(slide, "WHY US")

tb = add_textbox(slide, Inches(0.8), Inches(0.8), Inches(11.5), Inches(0.7))
set_text(tb.text_frame, "What We Understand That Others Don't", size=32, color=WHITE, bold=True)

insights = [
    (ACCENT, "The AI debate is not one discourse — it's three incompatible systems.",
     "Competitors build single-perspective tools. We build the translation layer between worldviews. This is the difference between a dictionary and a Rosetta Stone."),
    (ACCENT2, "The interesting information lives in the gaps between perspectives.",
     "Conflicts aren't noise — they're signal. Every disagreement, properly classified (values vs. facts vs. method), reveals exactly where policy can and cannot find common ground."),
    (ACCENT3, "Every new document makes every existing analysis richer.",
     "The structured taxonomy + conflict detection pipeline is a compounding data asset. Network effects: more sources ingested = more conflict surface area = higher diagnostic value for every user."),
    (RGBColor(0xBB, 0x86, 0xFC), "This is a dilemma, not a problem.",
     "Problems have solutions. Dilemmas have trade-offs. The goal is not to resolve the AI debate but to build the infrastructure that lets stakeholders make informed trade-offs."),
]

for i, (color, headline, detail) in enumerate(insights):
    top = Inches(1.7 + i * 1.35)
    add_shape_fill(slide, Inches(0.8), top, Inches(11.7), Inches(1.2), BG_LIGHT)
    add_shape_fill(slide, Inches(0.8), top, Inches(0.08), Inches(1.2), color)

    tb = add_textbox(slide, Inches(1.2), top + Inches(0.1), Inches(10.8), Inches(0.4))
    set_text(tb.text_frame, headline, size=16, color=color, bold=True)

    tb = add_textbox(slide, Inches(1.2), top + Inches(0.55), Inches(10.8), Inches(0.6))
    set_text(tb.text_frame, detail, size=13, color=LIGHT_GRAY)


# ============================================================
# SLIDE 10 — TEAM
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
section_label(slide, "WHY US")

tb = add_textbox(slide, Inches(0.8), Inches(0.8), Inches(11.5), Inches(0.7))
set_text(tb.text_frame, "Team: 40 Years of Building Bridges From Chaos", size=32, color=WHITE, bold=True)

# Name and headline
tb = add_textbox(slide, Inches(0.8), Inches(1.6), Inches(5), Inches(0.5))
set_text(tb.text_frame, "Jeffrey Snover", size=24, color=WHITE, bold=True)

tb = add_textbox(slide, Inches(0.8), Inches(2.1), Inches(5.5), Inches(0.5))
set_text(tb.text_frame, "Systems Philosopher  |  Inventor of PowerShell  |  Berkman Klein Center Fellow", size=13, color=ACCENT)

# Track record cards
records = [
    ("Unified Siloed Cultures Into Shared Taxonomies", ACCENT,
     "Co-authored \"Failure Modes in ML\" with Berkman Klein Center — first paper\n"
     "to jointly tabulate security attacks and safety bias. Became the foundation\nfor the MITRE ATLAS Framework (industry standard)."),
    ("Built Governance Frameworks From Chaos", ACCENT3,
     "At Google: synthesized divergent SRE practices across all services into a single\n"
     "Risk Taxonomy and Risk Control Catalog. Turned Tower of Babel into\nstandardized operational culture. Enabled scalable AI Ops tooling."),
    ("Drove Decade-Long Institutional Change", ACCENT2,
     "Invented PowerShell by writing the Monad Manifesto — reframed GUI vs. CLI\n"
     "as AND, not OR. Overcame deep cultural resistance to democratize\nautomation for millions of users worldwide."),
    ("Bridged Technical and Non-Technical Stakeholders", RGBColor(0xBB, 0x86, 0xFC),
     "CTO/AI Architect for M365. Company spokesperson. Key mediator in\n"
     "Microsoft's corporate-wide Open Source transition. Led Azure Stack strategy\ndelivering \"cloud on your terms.\""),
]

for i, (title, color, desc) in enumerate(records):
    col = i % 2
    row = i // 2
    left = Inches(0.8 + col * 6.2)
    top = Inches(2.7 + row * 1.9)
    add_shape_fill(slide, left, top, Inches(5.8), Inches(1.7), BG_LIGHT)
    add_shape_fill(slide, left, top, Inches(5.8), Inches(0.05), color)

    tb = add_textbox(slide, left + Inches(0.3), top + Inches(0.15), Inches(5.2), Inches(0.4))
    set_text(tb.text_frame, title, size=14, color=color, bold=True)

    tb = add_textbox(slide, left + Inches(0.3), top + Inches(0.55), Inches(5.2), Inches(1.1))
    set_text(tb.text_frame, desc, size=11, color=LIGHT_GRAY)

# Bottom stats bar
add_shape_fill(slide, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.7), BG_MED)

stats = ["35+ Patents", "USENIX LISA Award", "65K+ Community", "Founding Member,\nDevOps Enterprise Forum",
         "Google DE\n(2022-2026)", "MS Technical Fellow\n(1999-2022)"]
for i, stat in enumerate(stats):
    left = Inches(0.8 + i * 1.95)
    tb = add_textbox(slide, left + Inches(0.1), Inches(6.52), Inches(1.75), Inches(0.65))
    set_text(tb.text_frame, stat, size=10, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 11 — THE ASK
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
section_label(slide, "THE ASK")

tb = add_textbox(slide, Inches(0.8), Inches(0.8), Inches(11.5), Inches(0.7))
set_text(tb.text_frame, "The Ask", size=32, color=WHITE, bold=True)

tb = add_textbox(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(0.5))
set_text(tb.text_frame, "We are seeking design partners and domain collaborators to shape what comes next.",
         size=18, color=GRAY)

asks = [
    ("DESIGN PARTNERS", ACCENT,
     "3-5 policy teams willing to pilot",
     "Think tank analysts, congressional research staff, corporate AI governance leads.\n"
     "90-day pilot in exchange for direct roadmap influence.\n"
     "We need real workflows to validate the tool against."),
    ("DOMAIN COLLABORATORS", ACCENT3,
     "Journalists and academics to stress-test",
     "Journalists covering AI who will test the \"decoder ring\" against real reporting.\n"
     "Academics who will pressure-test the taxonomy against research workflows.\n"
     "Goal: ensure the framework is grounded outside the engineering bubble."),
    ("GUIDANCE", ACCENT2,
     "Three open questions we need help answering",
     "1.  Which output format is most valuable: desktop tool, web app, or API?\n"
     "2.  What policy domains beyond AI should we target next?\n"
     "3.  What's the right model: institutional license, open-access, or hybrid?"),
]

for i, (label, color, headline, desc) in enumerate(asks):
    top = Inches(2.3 + i * 1.6)
    add_shape_fill(slide, Inches(0.8), top, Inches(11.7), Inches(1.4), BG_LIGHT)
    add_shape_fill(slide, Inches(0.8), top, Inches(0.08), Inches(1.4), color)

    tb = add_textbox(slide, Inches(1.2), top + Inches(0.1), Inches(2), Inches(0.3))
    set_text(tb.text_frame, label, size=11, color=color, bold=True)

    tb = add_textbox(slide, Inches(1.2), top + Inches(0.4), Inches(4), Inches(0.4))
    set_text(tb.text_frame, headline, size=16, color=WHITE, bold=True)

    tb = add_textbox(slide, Inches(5.5), top + Inches(0.15), Inches(6.5), Inches(1.1))
    set_text(tb.text_frame, desc, size=12, color=LIGHT_GRAY)

# Closing quote
add_shape_fill(slide, Inches(2), Inches(6.3), Inches(9.3), Inches(0.8), BG_MED)
tb = add_textbox(slide, Inches(2.3), Inches(6.4), Inches(8.7), Inches(0.6))
set_text(tb.text_frame,
         "\"We are not engineering the solution; we are engineering the table\nat which the solution can be negotiated.\"",
         size=15, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

# Bottom accent bar
add_shape_fill(slide, Inches(0), Inches(7.44), SLIDE_W, Inches(0.06), ACCENT)


# ============================================================
# SAVE MAIN DECK
# ============================================================
out_path = "/Users/jsnover/source/repos/ai-triad-research/AI_Triad_Pitch_Deck.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")


# ============================================================
# BACKUP DECK — AIF EXPLAINED
# ============================================================
backup = Presentation()
backup.slide_width  = Inches(13.333)
backup.slide_height = Inches(7.5)

slide = backup.slides.add_slide(backup.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
section_label(slide, "BACKUP")

tb = add_textbox(slide, Inches(0.8), Inches(0.8), Inches(11.5), Inches(0.7))
set_text(tb.text_frame, "What Is AIF? (Argumentation Interchange Format)", size=32, color=WHITE, bold=True)

# Definition box
add_shape_fill(slide, Inches(0.8), Inches(1.7), Inches(11.7), Inches(1.3), BG_MED)
tb = add_textbox(slide, Inches(1.2), Inches(1.8), Inches(11), Inches(1.1))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "AIF is an international standard (developed at the University of Dundee) for representing arguments as structured data.",
         size=16, color=WHITE, bold=True)
add_para(tf, "", size=6, color=WHITE)
add_para(tf,
         "It provides a formal ontology for encoding claims, premises, inferences, and attacks — making arguments machine-readable "
         "and comparable across documents, authors, and viewpoints. Think of it as a graph schema for reasoning.",
         size=14, color=LIGHT_GRAY)

# Why it matters
tb = add_textbox(slide, Inches(0.8), Inches(3.2), Inches(5.5), Inches(0.5))
set_text(tb.text_frame, "Why It Matters for AI Triad", size=20, color=ACCENT, bold=True)

why_items = [
    ("Machine-comparable arguments", "Without AIF, summaries are free text. With AIF, every claim has typed relationships to other claims — we can algorithmically detect conflicts, not just keyword-match."),
    ("Attack classification", "AIF distinguishes rebuttals (\"your conclusion is wrong\"), undercuts (\"your reasoning is flawed\"), and undermines (\"your evidence is weak\"). This tells a policymaker what kind of disagreement they're looking at."),
    ("Cross-document linking", "Claims from different documents can be linked into a single argument graph. A Safetyist claim and an Accelerationist counter-claim become nodes in the same structure — visible, navigable, auditable."),
]

for i, (title, desc) in enumerate(why_items):
    top = Inches(3.8 + i * 1.15)
    add_shape_fill(slide, Inches(0.8), top, Inches(5.5), Inches(1.0), BG_LIGHT)
    add_shape_fill(slide, Inches(0.8), top, Inches(0.08), Inches(1.0), ACCENT)

    tb = add_textbox(slide, Inches(1.2), top + Inches(0.08), Inches(5), Inches(0.3))
    set_text(tb.text_frame, title, size=13, color=ACCENT, bold=True)

    tb = add_textbox(slide, Inches(1.2), top + Inches(0.38), Inches(5), Inches(0.6))
    set_text(tb.text_frame, desc, size=11, color=LIGHT_GRAY)

# Right side: edge types diagram
tb = add_textbox(slide, Inches(6.8), Inches(3.2), Inches(5.5), Inches(0.5))
set_text(tb.text_frame, "AI Triad's 7 AIF-Aligned Edge Types", size=20, color=ACCENT3, bold=True)

edge_types = [
    ("SUPPORTS", ACCENT, "Claim A provides evidence for Claim B"),
    ("CONTRADICTS", ACCENT2, "Claim A directly negates Claim B"),
    ("ASSUMES", LIGHT_GRAY, "Claim A depends on unstated Claim B"),
    ("WEAKENS", RGBColor(0xFF, 0xA5, 0x00), "Claim A reduces confidence in Claim B"),
    ("RESPONDS_TO", ACCENT3, "Claim A is a direct reply to Claim B"),
    ("TENSION_WITH", RGBColor(0xBB, 0x86, 0xFC), "Claims A and B are in unresolved friction"),
    ("INTERPRETS", RGBColor(0xFF, 0xD7, 0x00), "Claim A reframes the meaning of Claim B"),
]

for i, (etype, color, desc) in enumerate(edge_types):
    top = Inches(3.8 + i * 0.48)
    add_shape_fill(slide, Inches(6.8), top, Inches(5.5), Inches(0.42), BG_LIGHT)
    add_shape_fill(slide, Inches(6.8), top, Inches(0.08), Inches(0.42), color)

    tb = add_textbox(slide, Inches(7.1), top + Inches(0.04), Inches(2), Inches(0.35))
    set_text(tb.text_frame, etype, size=11, color=color, bold=True)

    tb = add_textbox(slide, Inches(9.1), top + Inches(0.04), Inches(3), Inches(0.35))
    set_text(tb.text_frame, desc, size=10, color=LIGHT_GRAY)

# Bottom note
tb = add_textbox(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.5))
set_text(tb.text_frame,
         "Reference: Chesnevar et al., \"Towards an Argument Interchange Format,\" The Knowledge Engineering Review (2006)",
         size=11, color=GRAY, alignment=PP_ALIGN.CENTER)

backup_path = "/Users/jsnover/source/repos/ai-triad-research/AI_Triad_Backup_Slides.pptx"
backup.save(backup_path)
print(f"Saved: {backup_path}")
