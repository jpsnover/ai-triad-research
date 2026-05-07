#!/usr/bin/env python3
"""Generate PowerPoint deck: Applying Computational Dialectics to Understand the AI Triad"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xD6)
ACCENT_GOLD = RGBColor(0xE8, 0xB0, 0x4D)
ACCENT_RED = RGBColor(0xE0, 0x5A, 0x5A)
ACCENT_GREEN = RGBColor(0x5A, 0xB8, 0x7A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MED_GRAY = RGBColor(0x99, 0x99, 0x99)
DARK_TEXT = RGBColor(0x2D, 0x2D, 0x2D)
SOFT_BG = RGBColor(0xF5, 0xF5, 0xF0)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Calibri", line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(font_size * 0.3)
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    return tf


def add_para(tf, text, font_size=18, color=DARK_TEXT, bold=False,
             alignment=PP_ALIGN.LEFT, font_name="Calibri", space_after=None,
             level=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.level = level
    p.space_after = Pt(space_after if space_after else font_size * 0.4)
    return p


def add_bullet_slide(slide, left, top, width, height, bullets, font_size=20,
                     color=DARK_TEXT, bullet_color=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(font_size * 0.6)
        p.level = 0
    return tf


def add_hyperlink_para(tf, text, url, font_size=16, color=ACCENT_BLUE,
                       font_name="Calibri", bold=False, space_after=None):
    """Add a paragraph with a clickable hyperlink."""
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name
    run.font.underline = True
    run.hyperlink.address = url
    p.space_after = Pt(space_after if space_after else font_size * 0.5)
    return p


def add_accent_bar(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_section_divider(slide, act_label, title, color=ACCENT_BLUE):
    set_slide_bg(slide, DARK_BG)
    add_accent_bar(slide, Inches(0), Inches(3.2), W, Inches(0.06), color)
    add_textbox(slide, Inches(1), Inches(2.0), Inches(11), Inches(0.8),
                act_label, font_size=22, color=color, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(3.6), Inches(11), Inches(1.2),
                title, font_size=40, color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER, font_name="Calibri Light")


# ============================================================
# SLIDE 1: TITLE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, DARK_BG)

add_accent_bar(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT_BLUE)

add_textbox(slide, Inches(1.2), Inches(1.5), Inches(11), Inches(1.5),
            "Applying Computational Dialectics\nto Understand the AI Triad",
            font_size=44, color=WHITE, bold=True, font_name="Calibri Light",
            alignment=PP_ALIGN.LEFT, line_spacing=1.3)

add_accent_bar(slide, Inches(1.2), Inches(3.5), Inches(3), Inches(0.05), ACCENT_GOLD)

add_textbox(slide, Inches(1.2), Inches(3.8), Inches(11), Inches(0.6),
            "Jeffrey Snover  |  Berkman Klein Center Fellow  |  2026",
            font_size=22, color=LIGHT_GRAY, alignment=PP_ALIGN.LEFT)

add_textbox(slide, Inches(1.2), Inches(5.2), Inches(10), Inches(1.5),
            "Using AI to systematically interrogate the AI policy debate —\n"
            "decomposing positions, staging structured debates, and mapping\n"
            "where camps agree, disagree, and what kind of disagreement it is.",
            font_size=18, color=MED_GRAY, alignment=PP_ALIGN.LEFT, line_spacing=1.4)

# ============================================================
# ACT I DIVIDER
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_divider(slide, "ACT I", "The Problem", ACCENT_RED)

# ============================================================
# SLIDE 2: THE AI TRIAD
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SOFT_BG)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "The AI Triad — Three Silo'd Monologues",
            font_size=36, color=DARK_TEXT, bold=True, font_name="Calibri Light")

add_accent_bar(slide, Inches(0.8), Inches(1.15), Inches(2.5), Inches(0.04), ACCENT_BLUE)

add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.8),
            "Different camps use the same words to mean different things,\n"
            "and different words to mean the same thing.",
            font_size=20, color=MED_GRAY)

# Three POV boxes
box_top = Inches(2.8)
box_h = Inches(3.5)
box_w = Inches(3.5)
gap = Inches(0.4)
start_left = Inches(0.8)

pov_data = [
    ("Accelerationists", "Prometheus", ACCENT_BLUE,
     ["AI capability growth is urgent",
      "Regulation is a bottleneck",
      "Competitive dynamics between\ndemocracies drive responsible\ndeployment"]),
    ("Safetyists", "Sentinel", ACCENT_GOLD,
     ["Current safety measures are\ninadequate",
      "Precautionary approaches are\njustified",
      "Institutional governance must\ngate high-risk capabilities"]),
    ("Skeptics", "Cassandra", ACCENT_RED,
     ["Hypothetical risks distract from\ndemonstrated harms",
      "Labor displacement, algorithmic\nbias, power concentration",
      "Accountability gaps demand\nimmediate action"]),
]

for i, (name, char, color, points) in enumerate(pov_data):
    left = start_left + i * (box_w + gap)
    # Card background
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, box_top, box_w, box_h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(1)
    # Color bar at top
    add_accent_bar(slide, left, box_top, box_w, Inches(0.06), color)
    # Title
    add_textbox(slide, left + Inches(0.2), box_top + Inches(0.2), box_w - Inches(0.4), Inches(0.5),
                name, font_size=22, color=color, bold=True)
    add_textbox(slide, left + Inches(0.2), box_top + Inches(0.65), box_w - Inches(0.4), Inches(0.3),
                f"({char})", font_size=16, color=MED_GRAY)
    # Bullets
    for j, pt in enumerate(points):
        add_textbox(slide, left + Inches(0.3), box_top + Inches(1.1 + j * 0.75),
                    box_w - Inches(0.5), Inches(0.7),
                    f"• {pt}", font_size=14, color=DARK_TEXT)

# Bottom note
add_textbox(slide, Inches(0.8), Inches(6.6), Inches(11), Inches(0.5),
            "Each has serious intellectual backing, coherent internal logic, and real policy influence.\n"
            "But they rarely engage each other's strongest arguments.",
            font_size=16, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 3: WHY THIS MATTERS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SOFT_BG)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "Why This Matters — The Targets of Influence",
            font_size=36, color=DARK_TEXT, bold=True, font_name="Calibri Light")
add_accent_bar(slide, Inches(0.8), Inches(1.15), Inches(2.5), Inches(0.04), ACCENT_BLUE)

# Three target groups
targets = [
    ("Policymakers", "Must write rules based on\ncontested facts and competing values", ACCENT_BLUE),
    ("Journalists", "Must explain a debate where\nthe vocabulary itself is weaponized", ACCENT_GOLD),
    ("Academics", "Must synthesize across traditions\nthat refuse to share terminology", ACCENT_RED),
]

for i, (role, desc, color) in enumerate(targets):
    top = Inches(1.8 + i * 1.2)
    add_accent_bar(slide, Inches(0.8), top, Inches(0.06), Inches(0.9), color)
    add_textbox(slide, Inches(1.1), top, Inches(3), Inches(0.5),
                role, font_size=24, color=color, bold=True)
    add_textbox(slide, Inches(1.1), top + Inches(0.4), Inches(5), Inches(0.5),
                desc, font_size=17, color=DARK_TEXT)

# Disagreement types
add_textbox(slide, Inches(7), Inches(1.8), Inches(5.5), Inches(0.6),
            "What type of disagreement is it?",
            font_size=24, color=DARK_TEXT, bold=True)

types = [
    ("Definitional", "Resolved with clarity of language", ACCENT_BLUE),
    ("Empirical", "Resolved by gathering critical data", ACCENT_GOLD),
    ("Normative", "Conflict over values — tradeoffs needed", ACCENT_RED),
]

for i, (t, desc, color) in enumerate(types):
    top = Inches(2.7 + i * 1.0)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(7), top, Inches(5.3), Inches(0.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = color
    shape.line.width = Pt(2)
    add_textbox(slide, Inches(7.2), top + Inches(0.05), Inches(2), Inches(0.35),
                t, font_size=18, color=color, bold=True)
    add_textbox(slide, Inches(7.2), top + Inches(0.38), Inches(4.8), Inches(0.35),
                desc, font_size=15, color=MED_GRAY)

# AI Rosetta Stone
add_textbox(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(1.2),
            'They needed a psycho-technology — a way of processing information that enables\n'
            'them to comprehend what is truly being said.\n'
            'I called it the AI Rosetta Stone.',
            font_size=20, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

# ============================================================
# ACT II DIVIDER
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_divider(slide, "ACT II", "The Journey", ACCENT_GOLD)

# ============================================================
# SLIDE 4: STARTING WITH WHAT I KNEW
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SOFT_BG)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "Starting with What I Knew — Taxonomy Building",
            font_size=36, color=DARK_TEXT, bold=True, font_name="Calibri Light")
add_accent_bar(slide, Inches(0.8), Inches(1.15), Inches(2.5), Inches(0.04), ACCENT_GOLD)

tf = add_textbox(slide, Inches(0.8), Inches(1.7), Inches(11), Inches(5),
                 "I started the way I've done lots of similar projects in the corporate world.",
                 font_size=22, color=DARK_TEXT)

add_para(tf, "", font_size=12, color=DARK_TEXT)
add_para(tf, "Read a lot. Put the information into Excel. Stare at it.\nCluster it. Combine things. Split things.",
         font_size=22, color=MED_GRAY)
add_para(tf, "", font_size=12, color=DARK_TEXT)
add_para(tf, "I had done this successfully when I generated Google's Technology Risk Taxonomy.",
         font_size=22, color=DARK_TEXT)
add_para(tf, "", font_size=12, color=DARK_TEXT)
add_para(tf, "At the end stages of that project, LLM tools were starting to become useful.\n"
         "I tried them for language consistency, tone, level of detail.",
         font_size=22, color=DARK_TEXT)
add_para(tf, "", font_size=12, color=DARK_TEXT)
add_para(tf, "But at some point I had the wild idea to ask it what I was missing.",
         font_size=22, color=DARK_TEXT, bold=True)
add_para(tf, "", font_size=12, color=DARK_TEXT)
add_para(tf, "It generated 10 missing elements. 7 were crap — but 3 were gold.",
         font_size=24, color=ACCENT_GOLD, bold=True)
add_para(tf, "", font_size=16, color=DARK_TEXT)
add_para(tf, "Can we use AI to understand AI?",
         font_size=28, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 5: COMPUTATIONAL LINGUISTICS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SOFT_BG)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "Down the Rabbit Hole — Computational Linguistics",
            font_size=36, color=DARK_TEXT, bold=True, font_name="Calibri Light")
add_accent_bar(slide, Inches(0.8), Inches(1.15), Inches(2.5), Inches(0.04), ACCENT_GOLD)

# Left column - the approach
add_textbox(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.5),
            "The Approach", font_size=22, color=ACCENT_BLUE, bold=True)

steps = [
    "1. Generate a seed taxonomy for each POV camp",
    "2. Ingest source documents using LLM prompts\n   to expand the taxonomy",
    "3. Refine iteratively — prompt quality drives\n   output quality",
]
for i, step in enumerate(steps):
    add_textbox(slide, Inches(1.0), Inches(2.2 + i * 1.0), Inches(5), Inches(0.9),
                step, font_size=18, color=DARK_TEXT)

# Right column - key techniques
add_textbox(slide, Inches(7), Inches(1.6), Inches(5.5), Inches(0.5),
            "Key Techniques Adopted", font_size=22, color=ACCENT_GOLD, bold=True)

# BDI box
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(7), Inches(2.2), Inches(5.5), Inches(2.0))
shape.fill.solid()
shape.fill.fore_color.rgb = WHITE
shape.line.color.rgb = ACCENT_GOLD
shape.line.width = Pt(1.5)

add_textbox(slide, Inches(7.2), Inches(2.3), Inches(5), Inches(0.4),
            "BDI Framework", font_size=20, color=ACCENT_GOLD, bold=True)
bdi_items = [
    "Beliefs — things people claim to be true",
    "Desires — what they want to be true",
    "Intentions — methods to make them true",
]
for i, item in enumerate(bdi_items):
    add_textbox(slide, Inches(7.4), Inches(2.8 + i * 0.4), Inches(5), Inches(0.4),
                f"• {item}", font_size=16, color=DARK_TEXT)

# DOLCE box
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(7), Inches(4.5), Inches(5.5), Inches(1.8))
shape.fill.solid()
shape.fill.fore_color.rgb = WHITE
shape.line.color.rgb = ACCENT_BLUE
shape.line.width = Pt(1.5)

add_textbox(slide, Inches(7.2), Inches(4.6), Inches(5), Inches(0.4),
            "Vocabulary over Structure", font_size=20, color=ACCENT_BLUE, bold=True)
add_textbox(slide, Inches(7.2), Inches(5.1), Inches(5), Inches(1.0),
            "OWL/RDF would limit LLM leverage.\n"
            "Instead: DOLCE genus-differentia definitions —\n"
            "friendly to both ontological analysis\n"
            "and LLM processing.",
            font_size=16, color=DARK_TEXT)

# Result
add_textbox(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.6),
            "Result: 565 taxonomy nodes across 3 POVs + cross-cutting situations, "
            "each BDI-categorized with precise definitions.",
            font_size=18, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 6: GENERATOR AND SELECTOR
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SOFT_BG)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "The Insight — Generator and Selector Functions",
            font_size=36, color=DARK_TEXT, bold=True, font_name="Calibri Light")
add_accent_bar(slide, Inches(0.8), Inches(1.15), Inches(2.5), Inches(0.04), ACCENT_GOLD)

add_textbox(slide, Inches(0.8), Inches(1.7), Inches(11), Inches(1.2),
            '"We\'re Harvard Law. We believe truth emerges through argumentation."',
            font_size=26, color=MED_GRAY, alignment=PP_ALIGN.CENTER,
            font_name="Calibri Light")
add_textbox(slide, Inches(0.8), Inches(2.6), Inches(11), Inches(0.5),
            "— Harvard Professor, Berkman Klein Center interview",
            font_size=16, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# Two boxes
for i, (title, desc, color, icon) in enumerate([
    ("Generator Function", "Computational linguistics —\nproducing candidate beliefs,\nexpanding the taxonomy,\ngrowing the knowledge base",
     ACCENT_GOLD, "→"),
    ("Selector Function", "Computational dialectics —\napplying evolutionary forces to\nstress-test belief systems,\ncull weak links, refine strong ones",
     ACCENT_BLUE, "⇒"),
]):
    left = Inches(1.5) + i * Inches(5.5)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, Inches(3.5), Inches(4.5), Inches(2.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = color
    shape.line.width = Pt(2)
    add_textbox(slide, left + Inches(0.3), Inches(3.7), Inches(3.8), Inches(0.5),
                title, font_size=24, color=color, bold=True)
    add_textbox(slide, left + Inches(0.3), Inches(4.3), Inches(3.8), Inches(1.5),
                desc, font_size=18, color=DARK_TEXT)

# Arrow between boxes
add_textbox(slide, Inches(6.0), Inches(4.3), Inches(1.3), Inches(0.8),
            "→", font_size=48, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.6),
            "A synthetic debate tool seemed perfect for this.\n"
            "That led me down the rabbit hole of computational dialectics.",
            font_size=20, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

# ============================================================
# ACT III DIVIDER
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_divider(slide, "ACT III", "What We Built", ACCENT_GREEN)

# ============================================================
# SLIDE 7: THE GENERATOR PIPELINE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SOFT_BG)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "The Generator — From Documents to Structured Knowledge",
            font_size=36, color=DARK_TEXT, bold=True, font_name="Calibri Light")
add_accent_bar(slide, Inches(0.8), Inches(1.15), Inches(2.5), Inches(0.04), ACCENT_GREEN)

pipeline_steps = [
    ("1  CONVERT", "173 source documents (PDFs, policy papers,\nblog posts, transcripts) normalized to Markdown", ACCENT_BLUE),
    ("2  EXTRACT (FIRE)", "Confidence-gated iterative extraction — claims\nscored for specificity, warrants, consistency.\nThe system knows what it doesn't know.", ACCENT_GOLD),
    ("3  MAP", "Embedding-based relevance scoring maps each\nclaim to taxonomy nodes. Unmapped concepts\nflagged as taxonomy gaps.", ACCENT_GREEN),
    ("4  STANDARDIZE", "35-term controlled vocabulary resolves the\n'same word, different meaning' problem.\nColloquial terms flagged for disambiguation.", ACCENT_RED),
    ("5  CLUSTER", "Agglomerative clustering (cosine distance)\ngroups related claims across POVs.", ACCENT_BLUE),
    ("6  EVOLVE", "When taxonomy changes, affected documents\nare re-summarized. Generator and taxonomy\nco-evolve.", ACCENT_GOLD),
]

for i, (step, desc, color) in enumerate(pipeline_steps):
    col = i % 3
    row = i // 3
    left = Inches(0.8) + col * Inches(4.1)
    top = Inches(1.6) + row * Inches(2.8)

    add_accent_bar(slide, left, top, Inches(0.06), Inches(2.2), color)
    add_textbox(slide, left + Inches(0.2), top + Inches(0.05), Inches(3.6), Inches(0.4),
                step, font_size=18, color=color, bold=True)
    add_textbox(slide, left + Inches(0.2), top + Inches(0.5), Inches(3.6), Inches(1.6),
                desc, font_size=15, color=DARK_TEXT)

# ============================================================
# SLIDE 8: CALIBRATION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SOFT_BG)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "Quality Through Calibration",
            font_size=36, color=DARK_TEXT, bold=True, font_name="Calibri Light")
add_accent_bar(slide, Inches(0.8), Inches(1.15), Inches(2.5), Inches(0.04), ACCENT_GREEN)

add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.6),
            "How do you know any of this is working? Every parameter is empirically calibrated.",
            font_size=22, color=MED_GRAY)

cal_items = [
    ("15 Parameters", "Extraction thresholds, relevance\ncutoffs, clustering distances,\nconvergence signals, phase\ntransitions — none hand-tuned", ACCENT_BLUE),
    ("Coverage Tracking", "What % of source claims are\naddressed in debates?\nTri-state: covered (>0.50),\npartial (>0.30), uncovered", ACCENT_GOLD),
    ("Threshold Optimization", "Quadratic fitting with confidence\ngates. Parameters only move\nwhen evidence is strong\n(n≥15 for high confidence)", ACCENT_GREEN),
]

for i, (title, desc, color) in enumerate(cal_items):
    left = Inches(0.8) + i * Inches(4.1)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, Inches(2.4), Inches(3.8), Inches(3.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = color
    shape.line.width = Pt(1.5)
    add_textbox(slide, left + Inches(0.2), Inches(2.6), Inches(3.4), Inches(0.5),
                title, font_size=22, color=color, bold=True)
    add_textbox(slide, left + Inches(0.2), Inches(3.2), Inches(3.4), Inches(2.0),
                desc, font_size=16, color=DARK_TEXT)

# Example callout
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(0.8), Inches(5.8), Inches(11.5), Inches(1.2))
shape.fill.solid()
shape.fill.fore_color.rgb = DARK_BG
shape.line.fill.background()
add_textbox(slide, Inches(1.2), Inches(5.95), Inches(10.8), Inches(0.8),
            "Example: Original relevance threshold of 0.30 admitted 93% of node pairs — effectively no filtering.\n"
            "Empirical calibration moved it to 0.45, admitting ~70%. Dramatically improved signal-to-noise.",
            font_size=17, color=WHITE, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 9: THE DEBATE ENGINE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SOFT_BG)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "The Selector — Three AI Agents, Structured Argumentation",
            font_size=36, color=DARK_TEXT, bold=True, font_name="Calibri Light")
add_accent_bar(slide, Inches(0.8), Inches(1.15), Inches(2.5), Inches(0.04), ACCENT_GREEN)

# Three characters
chars = [
    ("Prometheus", "Accelerationist", ACCENT_BLUE),
    ("Sentinel", "Safetyist", ACCENT_GOLD),
    ("Cassandra", "Skeptic", ACCENT_RED),
]
for i, (name, role, color) in enumerate(chars):
    left = Inches(0.8) + i * Inches(4.1)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, Inches(1.6), Inches(3.8), Inches(1.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = color
    shape.line.width = Pt(2)
    add_textbox(slide, left + Inches(0.2), Inches(1.7), Inches(3.4), Inches(0.4),
                name, font_size=22, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.2), Inches(2.1), Inches(3.4), Inches(0.3),
                role, font_size=16, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# Phases
add_textbox(slide, Inches(0.8), Inches(3.0), Inches(11), Inches(0.5),
            "Structured Phases — Not Open-Ended Chat",
            font_size=22, color=DARK_TEXT, bold=True)

phases = [
    ("Thesis-Antithesis", "Stake out positions,\nchallenge core claims", ACCENT_RED),
    ("Exploration", "Probe deeper, find cruxes,\ntest edge cases", ACCENT_GOLD),
    ("Synthesis", "Identify convergence, classify\nremaining disagreements", ACCENT_GREEN),
]
for i, (name, desc, color) in enumerate(phases):
    left = Inches(0.8) + i * Inches(4.1)
    add_accent_bar(slide, left, Inches(3.6), Inches(3.8), Inches(0.06), color)
    add_textbox(slide, left, Inches(3.8), Inches(3.8), Inches(0.4),
                name, font_size=20, color=color, bold=True)
    add_textbox(slide, left, Inches(4.3), Inches(3.8), Inches(0.8),
                desc, font_size=16, color=DARK_TEXT)

# 4-stage pipeline
add_textbox(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(0.5),
            "Each turn follows a 4-stage \"argue like a lawyer\" pipeline:",
            font_size=20, color=DARK_TEXT, bold=True)

stages = [
    ("BRIEF", "Summarize situation", "T=0.15"),
    ("PLAN", "Select moves & strategy", "T=0.4"),
    ("DRAFT", "Generate argument", "T=0.7"),
    ("CITE", "Map to taxonomy", "T=0.15"),
]
for i, (name, desc, temp) in enumerate(stages):
    left = Inches(0.8) + i * Inches(3.1)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, Inches(6.1), Inches(2.8), Inches(1.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BG
    shape.line.fill.background()
    add_textbox(slide, left + Inches(0.15), Inches(6.15), Inches(2.5), Inches(0.35),
                f"{name}  ({temp})", font_size=16, color=ACCENT_BLUE, bold=True)
    add_textbox(slide, left + Inches(0.15), Inches(6.5), Inches(2.5), Inches(0.4),
                desc, font_size=14, color=LIGHT_GRAY)

# ============================================================
# SLIDE 10: DIALECTICAL MOVES
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SOFT_BG)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "Dialectical Moves — Making Argumentation Rigorous",
            font_size=36, color=DARK_TEXT, bold=True, font_name="Calibri Light")
add_accent_bar(slide, Inches(0.8), Inches(1.15), Inches(2.5), Inches(0.04), ACCENT_GREEN)

add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
            "10 canonical moves grounded in argumentation theory (Walton, Pollock, Hamblin)",
            font_size=18, color=MED_GRAY)

moves = [
    ("DISTINGUISH", "Accept evidence, deny applicability", ACCENT_RED),
    ("COUNTEREXAMPLE", "Concrete case challenges a general claim", ACCENT_RED),
    ("REFRAME", "Shift the frame to reveal hidden structure", ACCENT_RED),
    ("EMPIRICAL CHALLENGE", "Dispute facts with counter-evidence", ACCENT_RED),
    ("UNDERCUT", "Attack the reasoning, not the conclusion", ACCENT_RED),
    ("BURDEN-SHIFT", "Challenge who bears burden of proof", ACCENT_RED),
    ("CONCEDE-AND-PIVOT", "Genuine concession + redirect", ACCENT_GREEN),
    ("EXTEND", "Build on another's point with new substance", ACCENT_GREEN),
    ("INTEGRATE", "Synthesize into novel position", ACCENT_GREEN),
    ("SPECIFY", "Force falsifiable predictions — name the crux", ACCENT_GOLD),
]

for i, (name, desc, color) in enumerate(moves):
    col = i % 2
    row = i // 2
    left = Inches(0.8) + col * Inches(6.2)
    top = Inches(2.1) + row * Inches(0.85)

    add_textbox(slide, left, top, Inches(2.8), Inches(0.4),
                name, font_size=16, color=color, bold=True)
    add_textbox(slide, left + Inches(2.8), top, Inches(3.2), Inches(0.4),
                desc, font_size=15, color=DARK_TEXT)
    # Subtle separator
    if row < 4:
        add_accent_bar(slide, left, top + Inches(0.65), Inches(5.8), Inches(0.01),
                       RGBColor(0xE0, 0xE0, 0xE0))

# Callout for SPECIFY
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.8))
shape.fill.solid()
shape.fill.fore_color.rgb = DARK_BG
shape.line.fill.background()
add_textbox(slide, Inches(1.2), Inches(6.5), Inches(10.8), Inches(0.6),
            "SPECIFY is critical — it's the only move that forces falsifiability into the open,\n"
            "requiring a debater to state what would change their mind.",
            font_size=17, color=ACCENT_GOLD, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 11: THE ACTIVE MODERATOR
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SOFT_BG)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "The Active Moderator — Monitoring and Intervening",
            font_size=36, color=DARK_TEXT, bold=True, font_name="Calibri Light")
add_accent_bar(slide, Inches(0.8), Inches(1.15), Inches(2.5), Inches(0.04), ACCENT_GREEN)

# Left: 7 signals
add_textbox(slide, Inches(0.8), Inches(1.5), Inches(6), Inches(0.5),
            "7 Deterministic Signals (no LLM calls)", font_size=20, color=ACCENT_BLUE, bold=True)

signals = [
    ("Move Disposition", "Confrontational vs. collaborative ratio"),
    ("Engagement Depth", "Citing each other, or talking past?"),
    ("Recycling Rate", "Arguments exhausted? (token overlap)"),
    ("Strongest Attack", "Best opposing argument — being ignored?"),
    ("Concession Opportunity", "Rigidity vs. sycophancy detection"),
    ("Position Delta", "Silent drift without conceding"),
    ("Crux Rate", "Are identified cruxes being engaged?"),
]

for i, (name, desc) in enumerate(signals):
    top = Inches(2.1 + i * 0.65)
    add_textbox(slide, Inches(0.8), top, Inches(2.8), Inches(0.35),
                name, font_size=15, color=ACCENT_BLUE, bold=True)
    add_textbox(slide, Inches(3.6), top, Inches(3.5), Inches(0.35),
                desc, font_size=14, color=DARK_TEXT)

# Right: interventions
add_textbox(slide, Inches(7.5), Inches(1.5), Inches(5), Inches(0.5),
            "Moderator Interventions", font_size=20, color=ACCENT_GOLD, bold=True)

interventions = [
    ("Elicitation", "Draw out unexplored positions"),
    ("Repair", "Fix failed turns, redirect tangents"),
    ("Reconciliation", "Surface hidden agreements"),
    ("Reflection", "Force summarizing what was learned"),
    ("Synthesis", "Guide toward classifying disagreements"),
    ("Gap Injection", "Strong unmade arguments from fresh AI"),
]

for i, (name, desc) in enumerate(interventions):
    top = Inches(2.1 + i * 0.75)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(7.5), top, Inches(5), Inches(0.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = ACCENT_GOLD
    shape.line.width = Pt(1)
    add_textbox(slide, Inches(7.7), top + Inches(0.02), Inches(2.2), Inches(0.3),
                name, font_size=15, color=ACCENT_GOLD, bold=True)
    add_textbox(slide, Inches(7.7), top + Inches(0.28), Inches(4.5), Inches(0.3),
                desc, font_size=13, color=DARK_TEXT)

# Bottom note
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(0.8), Inches(6.6), Inches(11.5), Inches(0.6))
shape.fill.solid()
shape.fill.fore_color.rgb = DARK_BG
shape.line.fill.background()
add_textbox(slide, Inches(1.2), Inches(6.65), Inches(10.8), Inches(0.5),
            "Recommendations are neural. Validation is deterministic. The moderator can't override the structure.",
            font_size=16, color=ACCENT_GOLD, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 12: QBAF
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SOFT_BG)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "Making Arguments Computable — QBAF",
            font_size=36, color=DARK_TEXT, bold=True, font_name="Calibri Light")
add_accent_bar(slide, Inches(0.8), Inches(1.15), Inches(2.5), Inches(0.04), ACCENT_GREEN)

# Formula
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(0.8), Inches(1.6), Inches(11.5), Inches(1.4))
shape.fill.solid()
shape.fill.fore_color.rgb = DARK_BG
shape.line.fill.background()

add_textbox(slide, Inches(1.2), Inches(1.65), Inches(10.8), Inches(0.4),
            "Quantified Bipolar Argumentation Frameworks (DF-QuAD, Rago et al. 2016)",
            font_size=16, color=LIGHT_GRAY)
add_textbox(slide, Inches(1.2), Inches(2.1), Inches(10.8), Inches(0.6),
            "strength(a) = base(a) × (1 - aggregated_attacks) × (1 + aggregated_support)",
            font_size=22, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER,
            font_name="Consolas")

# Key points
points = [
    ("Every claim has a base strength", "scored for evidential quality"),
    ("Attacks reduce strength", "rebut (1.0×), undercut (1.05×), undermine (1.1×)"),
    ("Supports amplify strength", "converging evidence strengthens positions"),
    ("Iterates until convergence", "with oscillation detection and damping"),
    ("Every outcome is traceable", "follow the math from conclusion to evidence"),
]

for i, (main, detail) in enumerate(points):
    top = Inches(3.3 + i * 0.7)
    add_textbox(slide, Inches(1.2), top, Inches(5), Inches(0.35),
                f"• {main}", font_size=18, color=DARK_TEXT, bold=True)
    add_textbox(slide, Inches(6.2), top, Inches(5.5), Inches(0.35),
                detail, font_size=16, color=MED_GRAY)

# Key finding
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.8))
shape.fill.solid()
shape.fill.fore_color.rgb = WHITE
shape.line.color.rgb = ACCENT_GOLD
shape.line.width = Pt(2)
add_textbox(slide, Inches(1.2), Inches(6.45), Inches(10.8), Inches(0.7),
            "Key finding: AI reliably scores Desires & Intentions (self-contained in text) but NOT Beliefs\n"
            "(require external verification). This is architectural, not a prompt failure → hybrid human-AI scoring.",
            font_size=16, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 13: WHAT THE SYSTEM REVEALS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SOFT_BG)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "What the System Reveals",
            font_size=36, color=DARK_TEXT, bold=True, font_name="Calibri Light")
add_accent_bar(slide, Inches(0.8), Inches(1.15), Inches(2.5), Inches(0.04), ACCENT_GREEN)

# Stats bar
add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.5),
            "94 debates  •  320 aggregated cruxes  •  565 taxonomy nodes  •  173 source documents",
            font_size=18, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER, bold=True)

reveals = [
    ("Cruxes", "The specific questions that, if answered, would change a position.\n"
     "Each classified: empirical (data), values (tradeoffs), definitional (vocabulary).",
     ACCENT_GOLD, "320 aggregated"),
    ("Hidden Convergence", "Positions that agree on policy actions despite completely\n"
     "different reasoning paths. Same destination, different routes.",
     ACCENT_GREEN, "Cross-POV"),
    ("Steelman Arguments", "The strongest version of each position, validated against\n"
     "opponents' actual commitments. Not strawmen.",
     ACCENT_BLUE, "Per debate"),
    ("Position Drift", "Detecting when participants silently shift positions\n"
     "without conceding. A sycophancy guard for honest discourse.",
     ACCENT_RED, "Embedding-based"),
]

for i, (title, desc, color, stat) in enumerate(reveals):
    col = i % 2
    row = i // 2
    left = Inches(0.8) + col * Inches(6.2)
    top = Inches(2.3) + row * Inches(2.3)

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, Inches(5.8), Inches(2.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = color
    shape.line.width = Pt(1.5)
    add_accent_bar(slide, left, top, Inches(0.06), Inches(2.0), color)
    add_textbox(slide, left + Inches(0.3), top + Inches(0.1), Inches(3.5), Inches(0.4),
                title, font_size=22, color=color, bold=True)
    add_textbox(slide, left + Inches(3.8), top + Inches(0.15), Inches(1.8), Inches(0.3),
                stat, font_size=13, color=MED_GRAY, alignment=PP_ALIGN.RIGHT)
    add_textbox(slide, left + Inches(0.3), top + Inches(0.6), Inches(5.2), Inches(1.2),
                desc, font_size=15, color=DARK_TEXT)

# ============================================================
# SLIDE 14: FAILURE MODES
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SOFT_BG)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "Guarding Against AI Failure Modes",
            font_size=36, color=DARK_TEXT, bold=True, font_name="Calibri Light")
add_accent_bar(slide, Inches(0.8), Inches(1.15), Inches(2.5), Inches(0.04), ACCENT_GREEN)

add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
            "Using AI to analyze AI debates requires honesty about failure modes.",
            font_size=20, color=MED_GRAY)

failures = [
    ("Hallucination", "Confidence-gated extraction\nwith verification loops (FIRE)", ACCENT_RED),
    ("Sycophancy", "Embedding-based drift detection\ncatches agents agreeing too readily", ACCENT_GOLD),
    ("Steelman\nFabrication", "Cross-encoder validation against\nopponent's actual commitments", ACCENT_BLUE),
    ("Missing\nArguments", "Fresh AI with no debate context\nidentifies unmade arguments", ACCENT_GREEN),
    ("Persona\nContamination", "Neutral evaluator assesses\nwith speaker identities stripped", MED_GRAY),
]

for i, (title, desc, color) in enumerate(failures):
    left = Inches(0.5) + i * Inches(2.5)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, Inches(2.2), Inches(2.3), Inches(3.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = color
    shape.line.width = Pt(1.5)
    add_accent_bar(slide, left, Inches(2.2), Inches(2.3), Inches(0.06), color)
    add_textbox(slide, left + Inches(0.15), Inches(2.4), Inches(2.0), Inches(0.7),
                title, font_size=16, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.15), Inches(3.2), Inches(2.0), Inches(1.5),
                desc, font_size=13, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

# Bottom principle
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(0.8), Inches(5.6), Inches(11.5), Inches(1.5))
shape.fill.solid()
shape.fill.fore_color.rgb = DARK_BG
shape.line.fill.background()
add_textbox(slide, Inches(1.2), Inches(5.7), Inches(10.8), Inches(0.5),
            "The Critical Design Principle",
            font_size=20, color=ACCENT_GOLD, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.2), Inches(6.2), Inches(10.8), Inches(0.7),
            "Every outcome is explained through deterministic graph traversal,\n"
            "not another neural judgment. When we say a position prevailed,\n"
            "we show exactly which arguments, attacks, and concessions led there.",
            font_size=17, color=WHITE, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 15: DEMO TEE-UP
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.0),
            "The Platform",
            font_size=48, color=WHITE, bold=True, font_name="Calibri Light",
            alignment=PP_ALIGN.CENTER)

add_accent_bar(slide, Inches(4.5), Inches(2.7), Inches(4), Inches(0.05), ACCENT_GREEN)

features = [
    "Browse 565 taxonomy nodes across all POVs with BDI categorization",
    "Explore conflicts with QBAF-computed argument strengths",
    "View 320 aggregated cruxes — filterable by type and resolution status",
    "Run debates and watch the argument network build in real time",
    "Trace from any conclusion back to its evidential chain",
]

for i, feat in enumerate(features):
    add_textbox(slide, Inches(2), Inches(3.2 + i * 0.65), Inches(9), Inches(0.5),
                f"→  {feat}", font_size=20, color=LIGHT_GRAY)

add_textbox(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.8),
            "Let me show you...",
            font_size=28, color=ACCENT_GREEN, bold=True, alignment=PP_ALIGN.CENTER)

# ============================================================
# ACT IV DIVIDER
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_divider(slide, "ACT IV", "What We Learned & The Ask", ACCENT_GOLD)

# ============================================================
# SLIDE 16: WHAT WE LEARNED
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SOFT_BG)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "What We Learned",
            font_size=36, color=DARK_TEXT, bold=True, font_name="Calibri Light")
add_accent_bar(slide, Inches(0.8), Inches(1.15), Inches(2.5), Inches(0.04), ACCENT_GOLD)

learnings = [
    ("Camps agree on more policy actions\nthan their rhetoric suggests",
     "Hidden convergence is real — and actionable", ACCENT_GREEN),
    ("Many 'disagreements' are actually definitional\n— the same concept named differently",
     "The vocabulary problem is solvable", ACCENT_BLUE),
    ("The hardest disagreements are values-based\n— and naming them as such is itself progress",
     "Classification enables targeted intervention", ACCENT_GOLD),
    ("AI is remarkably good at stress-testing positions\n— but humans must remain in the loop",
     "Hybrid human-AI scoring for empirical claims", ACCENT_RED),
]

for i, (finding, implication, color) in enumerate(learnings):
    top = Inches(1.7 + i * 1.3)
    add_accent_bar(slide, Inches(0.8), top, Inches(0.06), Inches(1.0), color)
    add_textbox(slide, Inches(1.1), top, Inches(7), Inches(0.9),
                finding, font_size=19, color=DARK_TEXT, bold=True)
    add_textbox(slide, Inches(8.5), top + Inches(0.15), Inches(4), Inches(0.7),
                implication, font_size=16, color=MED_GRAY)

add_textbox(slide, Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.8),
            "The AI Rosetta Stone isn't about replacing human deliberation.\n"
            "It's about equipping humans with a map of the argumentative landscape.",
            font_size=20, color=DARK_TEXT, alignment=PP_ALIGN.CENTER, bold=True)

# ============================================================
# SLIDE 17: BEYOND AI POLICY
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SOFT_BG)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "Beyond AI Policy",
            font_size=36, color=DARK_TEXT, bold=True, font_name="Calibri Light")
add_accent_bar(slide, Inches(0.8), Inches(1.15), Inches(2.5), Inches(0.04), ACCENT_GOLD)

add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.8),
            "Anywhere there are coherent camps talking past each other\n"
            "with different vocabularies and assumptions — this method can map the terrain.",
            font_size=22, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

domains = [
    ("Climate Policy", "Growth vs. Sustainability\nvs. Justice", ACCENT_GREEN),
    ("Bioethics", "Innovation vs. Precaution\nvs. Access", ACCENT_BLUE),
    ("Platform Governance", "Free Speech vs. Safety\nvs. Competition", ACCENT_GOLD),
    ("Public Health", "Individual Liberty vs.\nCollective Welfare vs. Equity", ACCENT_RED),
]

for i, (name, desc, color) in enumerate(domains):
    left = Inches(0.8) + i * Inches(3.1)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, Inches(2.8), Inches(2.8), Inches(2.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = color
    shape.line.width = Pt(2)
    add_accent_bar(slide, left, Inches(2.8), Inches(2.8), Inches(0.06), color)
    add_textbox(slide, left + Inches(0.2), Inches(3.0), Inches(2.4), Inches(0.5),
                name, font_size=20, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.2), Inches(3.6), Inches(2.4), Inches(1.2),
                desc, font_size=16, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 18: THE ASK
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_textbox(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8),
            "The Ask",
            font_size=44, color=WHITE, bold=True, font_name="Calibri Light",
            alignment=PP_ALIGN.CENTER)
add_accent_bar(slide, Inches(5), Inches(1.4), Inches(3), Inches(0.05), ACCENT_GOLD)

# Use it
add_textbox(slide, Inches(1), Inches(1.8), Inches(11), Inches(0.5),
            "Use it.", font_size=28, color=ACCENT_GREEN, bold=True)
add_textbox(slide, Inches(1), Inches(2.3), Inches(11), Inches(0.8),
            "The tool is live. The system adapts output for your audience — policymaker language,\n"
            "researcher language, public language. Tell me where it breaks. Where it surprises you.",
            font_size=18, color=LIGHT_GRAY)

# Connect me
add_textbox(slide, Inches(1), Inches(3.4), Inches(11), Inches(0.5),
            "Connect me.", font_size=28, color=ACCENT_GOLD, bold=True)

asks = [
    ("Policymakers", "Can crux classification help them write\nbetter questions in hearings?", ACCENT_BLUE),
    ("Journalists", "Can steelman features help them\nrepresent positions fairly?", ACCENT_GOLD),
    ("Academics", "Can the auditable pipeline meet\nevidentiary standards?", ACCENT_RED),
]

for i, (role, question, color) in enumerate(asks):
    left = Inches(1) + i * Inches(3.9)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, Inches(4.1), Inches(3.6), Inches(1.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x25, 0x25, 0x40)
    shape.line.color.rgb = color
    shape.line.width = Pt(1.5)
    add_textbox(slide, left + Inches(0.2), Inches(4.2), Inches(3.2), Inches(0.4),
                role, font_size=20, color=color, bold=True)
    add_textbox(slide, left + Inches(0.2), Inches(4.7), Inches(3.2), Inches(1.0),
                question, font_size=15, color=LIGHT_GRAY)

# Final line
add_textbox(slide, Inches(1), Inches(6.3), Inches(11), Inches(0.8),
            "If you know someone who should see this, please introduce us.",
            font_size=24, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# ============================================================
# APPENDIX DIVIDER
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_divider(slide, "APPENDIX", "References & Documentation", MED_GRAY)

# ============================================================
# APPENDIX A: METHODS AND ALGORITHMS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SOFT_BG)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
            "Appendix A: Methods and Algorithms",
            font_size=32, color=DARK_TEXT, bold=True, font_name="Calibri Light")
add_accent_bar(slide, Inches(0.8), Inches(0.95), Inches(2.5), Inches(0.04), ACCENT_BLUE)

# Methods - left column (items 1-5)
methods_left = [
    ("QBAF", "Quantitative Bipolar Argumentation\nFrameworks with DF-QuAD gradual\nsemantics and BDI-aware calibration",
     [("Rago et al. (2016) — DF-QuAD",
       "https://aaai.org/papers/8-12874-discontinuity-free-decision-support-with-quantitative-argumentation-debates/"),
      ("Dung (1995) — Abstract Argumentation",
       "https://www.sciencedirect.com/science/article/pii/000437029400041X")]),
    ("BDI Framework", "Belief-Desire-Intention agent\ncharacterization for multi-perspective\nstance decomposition",
     [("Rao & Georgeff (1995)",
       "https://cdn.aaai.org/ICMAS/1995/ICMAS95-042.pdf"),
      ("Bratman (1987) — Intention, Plans",
       "https://press.uchicago.edu/ucp/books/book/distributed/I/bo3629095.html")]),
    ("AIF", "Argument Interchange Format —\ntyped attack/support relationships",
     [("Chesñevar et al. (2006)",
       "https://dl.acm.org/doi/10.1017/S0269888906001044")]),
    ("FIRE Extraction", "Confidence-gated iterative extraction\nreplacing single-shot claim extraction",
     [("arXiv:2411.00784",
       "https://arxiv.org/abs/2411.00784")]),
    ("4-Stage Pipeline", "BRIEF → PLAN → DRAFT → CITE\nwith per-stage temperature control", []),
]

# Methods - right column (items 6-9)
methods_right = [
    ("Adaptive Staging", "Seven convergence diagnostics\ntrigger phase transitions\ndeterministically", []),
    ("13-Scheme Taxonomy", "Argumentation schemes derived\nfrom Walton's framework",
     [("Walton, Reed & Macagno (2008)",
       "https://www.cambridge.org/core/books/argumentation-schemes/9AE7E4E6ABDE690565442B2BD516A8B6")]),
    ("14-Move Moderator", "Six intervention families governed\nby pragma-dialectical theory",
     [("van Eemeren & Grootendorst (2004)",
       "https://www.cambridge.org/us/catalogue/catalogue.asp?isbn=9780521537728")]),
    ("DOLCE Ontology", "Descriptive Ontology for Linguistic\nand Cognitive Engineering —\nvocabulary over structure",
     [("Borgo et al. (2023)",
       "https://arxiv.org/pdf/2308.01597")]),
]

def render_method_column(slide, methods, start_left, start_top):
    for i, (name, desc, links) in enumerate(methods):
        top = start_top + i * Inches(1.18)
        add_textbox(slide, start_left, top, Inches(2.5), Inches(0.35),
                    name, font_size=16, color=ACCENT_BLUE, bold=True)
        add_textbox(slide, start_left, top + Inches(0.32), Inches(5.5), Inches(0.55),
                    desc, font_size=12, color=DARK_TEXT)
        if links:
            txBox = slide.shapes.add_textbox(start_left + Inches(0.1),
                                             top + Inches(0.75),
                                             Inches(5.4), Inches(0.4))
            tf = txBox.text_frame
            tf.word_wrap = True
            for j, (link_text, url) in enumerate(links):
                if j == 0:
                    p = tf.paragraphs[0]
                    run = p.add_run()
                else:
                    p = tf.add_paragraph()
                    run = p.add_run()
                run.text = f"→ {link_text}"
                run.font.size = Pt(10)
                run.font.color.rgb = ACCENT_BLUE
                run.font.name = "Calibri"
                run.font.underline = True
                run.hyperlink.address = url
                p.space_after = Pt(2)

render_method_column(slide, methods_left, Inches(0.8), Inches(1.2))
render_method_column(slide, methods_right, Inches(7.0), Inches(1.2))

# Vertical divider
add_accent_bar(slide, Inches(6.7), Inches(1.2), Inches(0.02), Inches(5.8), RGBColor(0xDD, 0xDD, 0xDD))

# ============================================================
# APPENDIX B: DOCUMENTATION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SOFT_BG)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
            "Appendix B: Documentation",
            font_size=32, color=DARK_TEXT, bold=True, font_name="Calibri Light")
add_accent_bar(slide, Inches(0.8), Inches(0.95), Inches(2.5), Inches(0.04), ACCENT_GOLD)

REPO_BASE = "https://github.com/jsnover/ai-triad-research/blob/main"

docs = [
    ("Architecture Overview", "docs/architecture-overview.md",
     "Two-repo split, Electron apps, AI backends, data model"),
    ("Debate Engine Design", "docs/debate-engine-design.md",
     "Three-agent BDI debate system, QBAF scoring, moderator"),
    ("Debate System Overview", "docs/debate-system-overview.md",
     "High-level overview of the multi-agent debate system"),
    ("Theory of Success", "docs/theory-of-success.md",
     "Success criteria, step-by-step execution, known weaknesses"),
    ("Taxonomy & Ontology Guide", "docs/taxonomy-ontology-guide.md",
     "BDI categories, POVs, situations, genus-differentia definitions"),
    ("Document Processing Pipeline", "docs/document-processing-pipeline.md",
     "Document ingestion, chunking, and claim extraction"),
    ("FIRE Extraction", "docs/fire-extraction.md",
     "Confidence-gated iterative claim extraction details"),
    ("Rhetorical Strategies", "docs/rhetorical-strategies.md",
     "Argumentation strategies employed by debate agents"),
    ("Epistemic Types", "docs/epistemic-types.md",
     "Types of knowledge claims and evaluation criteria"),
    ("Emotional Registers", "docs/emotional-registers.md",
     "Character voice and register design for debate personas"),
    ("Comp. Dialectics Comparison", "docs/computational-dialectics-comparison.md",
     "How this system compares to other argumentation tools"),
    ("Adaptive Debate Staging", "docs/design/adaptive-debate-staging.md",
     "Convergence diagnostics and phase transition design"),
    ("Full Methodology Paper", "docs/academic-paper-draft.md",
     "Complete technical paper with algorithms and evaluation"),
]

# Render in two columns
col_split = 7  # items in left column

for i, (title, path, desc) in enumerate(docs):
    col = 0 if i < col_split else 1
    row = i if i < col_split else i - col_split
    left = Inches(0.8) + col * Inches(6.2)
    top = Inches(1.3) + row * Inches(0.85)

    # Title with hyperlink
    txBox = slide.shapes.add_textbox(left, top, Inches(5.8), Inches(0.35))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(15)
    run.font.color.rgb = ACCENT_BLUE
    run.font.bold = True
    run.font.name = "Calibri"
    run.font.underline = True
    run.hyperlink.address = f"{REPO_BASE}/{path}"
    p.space_after = Pt(1)

    # Description
    add_textbox(slide, left + Inches(0.1), top + Inches(0.32), Inches(5.7), Inches(0.4),
                desc, font_size=12, color=MED_GRAY)

# Vertical divider
add_accent_bar(slide, Inches(6.7), Inches(1.3), Inches(0.02), Inches(5.5), RGBColor(0xDD, 0xDD, 0xDD))

# Footer
add_textbox(slide, Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.4),
            "All documentation is available in the project's Help dialog (? icon) within the Taxonomy Editor.",
            font_size=14, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# SAVE
# ============================================================
output_path = "/Users/jsnover/source/repos/ai-triad-research/management/project-manager/AI-Triad-Computational-Dialectics.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
print(f"Total slides: {len(prs.slides)}")
